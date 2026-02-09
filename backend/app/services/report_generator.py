"""Report generation service with charts and visualizations.

Generates audit reports in various formats:
- PowerPoint presentations with native editable charts
- PDF documents with embedded chart images
"""

import io
from dataclasses import dataclass
from datetime import datetime
from typing import Any

import matplotlib

matplotlib.use("Agg")

# isort: split

import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt
from reportlab.lib import colors as rl_colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import (
    Image,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

from app.db import DuckDBManager

# matplotlib 日本語フォント設定
plt.rcParams["font.family"] = ["Meiryo", "Yu Gothic", "IPAGothic", "sans-serif"]
plt.rcParams["axes.unicode_minus"] = False

# ベンフォードの法則 理論値
BENFORD_EXPECTED = {
    1: 0.301, 2: 0.176, 3: 0.125, 4: 0.097, 5: 0.079,
    6: 0.067, 7: 0.058, 8: 0.051, 9: 0.046,
}


@dataclass
class ReportConfig:
    """Report generation configuration."""

    fiscal_year: int
    period_from: int | None = None
    period_to: int | None = None
    company_name: str = "Sample Company"
    report_title: str = "仕訳検証レポート"
    prepared_by: str = "JAIA"
    include_details: bool = True


class PPTReportGenerator:
    """PowerPoint report generator with native charts and visual elements."""

    # カラースキーム - プロフェッショナルネイビー系
    PRIMARY = RGBColor(0x1A, 0x36, 0x5D)
    SECONDARY = RGBColor(0x4A, 0x5A, 0x6A)
    ACCENT = RGBColor(0x29, 0x80, 0xB9)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_BG = RGBColor(0xF0, 0xF2, 0xF5)

    # リスクレベル別カラー
    COLOR_HIGH = RGBColor(0xE7, 0x4C, 0x3C)
    COLOR_MEDIUM = RGBColor(0xF3, 0x9C, 0x12)
    COLOR_LOW = RGBColor(0x29, 0x80, 0xB9)
    COLOR_MINIMAL = RGBColor(0x27, 0xAE, 0x60)

    # フォント
    FONT_JP = "Meiryo"

    def __init__(self, config: ReportConfig):
        """Initialize PPT generator."""
        self.config = config
        self.db = DuckDBManager()
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

        # データキャッシュ
        self._stats: dict[str, Any] = {}
        self._risk_dist: dict[str, int] = {}
        self._findings: list[dict[str, Any]] = []
        self._benford_digits: list[dict[str, Any]] = []
        self._benford_summary: dict[str, Any] = {}
        self._monthly_trend: list[dict[str, Any]] = []

    def generate(self) -> bytes:
        """Generate PPT report with charts.

        Returns:
            PPT file as bytes.
        """
        self._prefetch_data()

        self._add_title_slide()
        self._add_summary_slide()
        self._add_metrics_slide()
        self._add_risk_chart_slide()
        self._add_findings_chart_slide()
        self._add_benford_chart_slide()
        self._add_trend_chart_slide()
        self._add_recommendations_slide()
        self._add_next_steps_slide()
        self._add_appendix_slide()

        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output.getvalue()

    # ========== データ取得 ==========

    def _prefetch_data(self) -> None:
        """Pre-fetch all data for report generation."""
        self._stats = self._query_summary_stats()
        self._risk_dist = self._query_risk_distribution()
        self._findings = self._query_top_findings()
        digits, summary = self._query_benford_data()
        self._benford_digits = digits
        self._benford_summary = summary
        self._monthly_trend = self._query_monthly_trend()

    def _query_summary_stats(self) -> dict[str, Any]:
        """Get summary statistics from DB."""
        query = """
            SELECT
                COUNT(*) as total_entries,
                COUNT(DISTINCT journal_id) as total_journals,
                COALESCE(SUM(ABS(amount)), 0) as total_amount,
                MIN(effective_date) as date_from,
                MAX(effective_date) as date_to,
                SUM(CASE WHEN risk_score >= 60 THEN 1 ELSE 0 END) as high_risk_count,
                AVG(CASE WHEN risk_score > 0 THEN risk_score END) as avg_risk_score
            FROM journal_entries
            WHERE fiscal_year = ?
        """
        result = self.db.execute(query, [self.config.fiscal_year])
        if result:
            row = result[0]
            total = row[0] or 1
            return {
                "total_entries": row[0] or 0,
                "total_journals": row[1] or 0,
                "total_amount": row[2] or 0,
                "date_from": str(row[3]) if row[3] else "",
                "date_to": str(row[4]) if row[4] else "",
                "high_risk_count": row[5] or 0,
                "high_risk_pct": (row[5] or 0) / total * 100,
                "avg_risk_score": row[6] or 0,
            }
        return {}

    def _query_risk_distribution(self) -> dict[str, int]:
        """Get risk distribution counts."""
        query = """
            SELECT
                CASE
                    WHEN risk_score >= 60 THEN 'high'
                    WHEN risk_score >= 40 THEN 'medium'
                    WHEN risk_score >= 20 THEN 'low'
                    ELSE 'minimal'
                END as level,
                COUNT(*) as count
            FROM journal_entries
            WHERE fiscal_year = ?
            GROUP BY level
        """
        result = self.db.execute(query, [self.config.fiscal_year])
        dist = {"high": 0, "medium": 0, "low": 0, "minimal": 0}
        for row in result:
            dist[row[0]] = row[1]
        return dist

    def _query_top_findings(self) -> list[dict[str, Any]]:
        """Get top rule violation findings."""
        query = """
            SELECT rule_name, COUNT(*) as count
            FROM rule_violations rv
            JOIN journal_entries je ON rv.gl_detail_id = je.gl_detail_id
            WHERE je.fiscal_year = ?
            GROUP BY rule_name
            ORDER BY count DESC
            LIMIT 10
        """
        result = self.db.execute(query, [self.config.fiscal_year])
        return [{"rule_name": row[0], "count": row[1]} for row in result]

    def _query_benford_data(self) -> tuple[list[dict], dict]:
        """Get Benford analysis with per-digit data."""
        query = """
            SELECT
                CAST(SUBSTR(CAST(ABS(CAST(amount AS BIGINT)) AS VARCHAR), 1, 1) AS INTEGER) as digit,
                COUNT(*) as count
            FROM journal_entries
            WHERE fiscal_year = ?
                AND ABS(amount) >= 10
            GROUP BY digit
            HAVING digit BETWEEN 1 AND 9
            ORDER BY digit
        """
        result = self.db.execute(query, [self.config.fiscal_year])

        total = sum(row[1] for row in result) or 1
        digits = []
        mad = 0
        for digit_num in range(1, 10):
            count = 0
            for row in result:
                if row[0] == digit_num:
                    count = row[1]
                    break
            actual = count / total
            expected = BENFORD_EXPECTED.get(digit_num, 0)
            mad += abs(actual - expected)
            digits.append({
                "digit": digit_num,
                "count": count,
                "actual_pct": round(actual * 100, 2),
                "expected_pct": round(expected * 100, 2),
            })
        mad = mad / 9

        if mad <= 0.006:
            label = "非常に良い適合"
            conformity = "close"
        elif mad <= 0.012:
            label = "許容範囲"
            conformity = "acceptable"
        elif mad <= 0.015:
            label = "境界"
            conformity = "marginally_acceptable"
        else:
            label = "不適合（要調査）"
            conformity = "nonconforming"

        summary = {"mad": mad, "conformity": conformity, "conformity_label": label}
        return digits, summary

    def _query_monthly_trend(self) -> list[dict[str, Any]]:
        """Get monthly trend data."""
        query = """
            SELECT
                accounting_period as period,
                COUNT(*) as entry_count,
                COALESCE(SUM(ABS(amount)), 0) as total_amount,
                SUM(CASE WHEN risk_score >= 60 THEN 1 ELSE 0 END) as high_risk_count
            FROM journal_entries
            WHERE fiscal_year = ?
            GROUP BY accounting_period
            ORDER BY accounting_period
        """
        result = self.db.execute(query, [self.config.fiscal_year])
        return [
            {
                "period": row[0],
                "entry_count": row[1],
                "total_amount": row[2],
                "high_risk_count": row[3],
            }
            for row in result
        ]

    # ========== ヘルパー ==========

    def _add_header_bar(self, slide: Any) -> None:
        """Add branded header bar at top of slide."""
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(0.08),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.PRIMARY
        bar.line.width = Emu(0)

    def _add_slide_title(self, slide: Any, title: str) -> None:
        """Add styled title to slide with header bar."""
        self._add_header_bar(slide)

        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(0.25), Inches(11.833), Inches(0.8),
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = self.PRIMARY
        p.font.name = self.FONT_JP

    def _add_slide_number(self, slide: Any, num: int, total: int = 10) -> None:
        """Add slide number at bottom right."""
        num_box = slide.shapes.add_textbox(
            Inches(12), Inches(7.0), Inches(1), Inches(0.35),
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{num} / {total}"
        p.font.size = Pt(10)
        p.font.color.rgb = self.SECONDARY
        p.alignment = PP_ALIGN.RIGHT

    def _add_kpi_card(
        self, slide: Any,
        x: float, y: float, w: float, h: float,
        label: str, value: str, accent_color: RGBColor,
    ) -> None:
        """Add a KPI card with colored top border."""
        # カード背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.WHITE
        card.line.color.rgb = RGBColor(0xDE, 0xE2, 0xE6)
        card.line.width = Pt(1)

        # カラーアクセントバー（上部）
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x + 0.05), Inches(y + 0.05),
            Inches(w - 0.1), Inches(0.06),
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_color
        accent_bar.line.width = Emu(0)

        # 数値
        val_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 0.35), Inches(w - 0.2), Inches(0.8),
        )
        tf = val_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.PRIMARY
        p.alignment = PP_ALIGN.CENTER

        # ラベル
        lbl_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 1.2), Inches(w - 0.2), Inches(0.5),
        )
        tf = lbl_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(13)
        p.font.color.rgb = self.SECONDARY
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.FONT_JP

    def _add_traffic_light(
        self, slide: Any, x: float, y: float, color: RGBColor,
    ) -> None:
        """Add a traffic light indicator (colored circle)."""
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(y), Inches(0.5), Inches(0.5),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.width = Emu(0)

    def _get_assessment(self) -> tuple[str, str, RGBColor]:
        """Get overall assessment text and color."""
        pct = self._stats.get("high_risk_pct", 0)
        if pct > 5:
            return (
                "要注意",
                "高リスク仕訳の割合が高く、詳細調査が必要です",
                self.COLOR_HIGH,
            )
        elif pct > 2:
            return (
                "注意",
                "一部注意を要する仕訳があります",
                self.COLOR_MEDIUM,
            )
        else:
            return (
                "概ね良好",
                "重大な問題は検出されていません",
                self.COLOR_MINIMAL,
            )

    # ========== スライド生成 ==========

    def _add_title_slide(self) -> None:
        """Slide 1: Title slide with branding."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 背景バー（装飾）
        bg_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(13.333), Inches(2.5),
        )
        bg_bar.fill.solid()
        bg_bar.fill.fore_color.rgb = self.PRIMARY
        bg_bar.line.width = Emu(0)

        # JAIA ロゴテキスト
        logo_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(0.5), Inches(4), Inches(0.6),
        )
        tf = logo_box.text_frame
        p = tf.paragraphs[0]
        p.text = "JAIA"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.WHITE

        sub = tf.add_paragraph()
        sub.text = "Journal entry AI Analyzer"
        sub.font.size = Pt(11)
        sub.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)

        # メインタイトル
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(3.2), Inches(11.833), Inches(1.2),
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = self.config.report_title
        p.font.size = Pt(42)
        p.font.bold = True
        p.font.color.rgb = self.PRIMARY
        p.font.name = self.FONT_JP

        # サブタイトル
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(4.5), Inches(11.833), Inches(0.8),
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{self.config.fiscal_year}年度  {self.config.company_name}"
        p.font.size = Pt(24)
        p.font.color.rgb = self.SECONDARY
        p.font.name = self.FONT_JP

        # 日付・作成者
        date_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(5.8), Inches(11.833), Inches(0.8),
        )
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        now = datetime.now()
        p.text = f"作成日: {now.strftime('%Y年%m月%d日')}　作成者: {self.config.prepared_by}"
        p.font.size = Pt(13)
        p.font.color.rgb = self.SECONDARY
        p.font.name = self.FONT_JP

    def _add_summary_slide(self) -> None:
        """Slide 2: Executive summary with traffic light."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, "エグゼクティブサマリー")
        self._add_slide_number(slide, 2)

        stats = self._stats
        assessment, detail, color = self._get_assessment()

        # 総合評価セクション（左側）
        self._add_traffic_light(slide, 0.75, 1.5, color)

        eval_box = slide.shapes.add_textbox(
            Inches(1.45), Inches(1.45), Inches(4), Inches(0.65),
        )
        tf = eval_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"総合評価: {assessment}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = self.FONT_JP

        detail_box = slide.shapes.add_textbox(
            Inches(1.45), Inches(2.1), Inches(4.5), Inches(0.5),
        )
        tf = detail_box.text_frame
        p = tf.paragraphs[0]
        p.text = detail
        p.font.size = Pt(14)
        p.font.color.rgb = self.SECONDARY
        p.font.name = self.FONT_JP

        # 検証概要（左側下）
        overview_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(3.0), Inches(5.5), Inches(4),
        )
        tf = overview_box.text_frame
        tf.word_wrap = True

        lines = [
            ("検証概要", True, Pt(16), self.PRIMARY),
            (f"検証対象仕訳件数: {stats.get('total_entries', 0):,}件", False, Pt(14), self.SECONDARY),
            (f"総取引金額: ¥{stats.get('total_amount', 0):,.0f}", False, Pt(14), self.SECONDARY),
            (f"対象期間: {stats.get('date_from', '')} ～ {stats.get('date_to', '')}", False, Pt(14), self.SECONDARY),
            ("", False, Pt(8), self.SECONDARY),
            ("リスク評価", True, Pt(16), self.PRIMARY),
            (f"高リスク仕訳: {stats.get('high_risk_count', 0):,}件 ({stats.get('high_risk_pct', 0):.1f}%)", False, Pt(14), self.SECONDARY),
            (f"平均リスクスコア: {stats.get('avg_risk_score', 0):.1f}点", False, Pt(14), self.SECONDARY),
        ]

        for i, (text, bold, size, clr) in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.font.size = size
            p.font.bold = bold
            p.font.color.rgb = clr
            p.font.name = self.FONT_JP
            p.space_before = Pt(4 if not bold else 12)

        # リスク分布ミニチャート（右側）- ドーナツ
        dist = self._risk_dist
        categories = ["高リスク", "中リスク", "低リスク", "最小リスク"]
        values = [dist.get("high", 0), dist.get("medium", 0),
                  dist.get("low", 0), dist.get("minimal", 0)]

        if sum(values) > 0:
            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series("件数", tuple(values))

            chart_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.DOUGHNUT,
                Inches(7), Inches(1.5), Inches(5.5), Inches(5),
                chart_data,
            )
            chart = chart_frame.chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False

            # ドーナツの色を設定
            plot = chart.plots[0]
            series = plot.series[0]
            risk_colors = [self.COLOR_HIGH, self.COLOR_MEDIUM,
                           self.COLOR_LOW, self.COLOR_MINIMAL]
            for i, clr in enumerate(risk_colors):
                point = series.points[i]
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = clr

    def _add_metrics_slide(self) -> None:
        """Slide 3: Key metrics dashboard with KPI cards."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, "主要指標ダッシュボード")
        self._add_slide_number(slide, 3)

        stats = self._stats

        # 上段: 4つのKPIカード
        cards = [
            ("検証仕訳件数", f"{stats.get('total_entries', 0):,}", self.ACCENT),
            ("仕訳帳票数", f"{stats.get('total_journals', 0):,}", self.ACCENT),
            ("高リスク件数", f"{stats.get('high_risk_count', 0):,}", self.COLOR_HIGH),
            ("平均リスクスコア", f"{stats.get('avg_risk_score', 0):.1f}", self.COLOR_MEDIUM),
        ]

        x_start = 0.75
        card_w = 2.8
        gap = 0.25
        for i, (label, value, color) in enumerate(cards):
            x = x_start + i * (card_w + gap)
            self._add_kpi_card(slide, x, 1.5, card_w, 1.9, label, value, color)

        # 下段: 追加指標テーブル
        info_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(4.0), Inches(11.833), Inches(3),
        )
        tf = info_box.text_frame
        tf.word_wrap = True

        extra_lines = [
            ("詳細指標", True, Pt(16), self.PRIMARY),
            (f"総取引金額: ¥{stats.get('total_amount', 0):,.0f}", False, Pt(14), self.SECONDARY),
            (f"対象期間: {stats.get('date_from', '')} ～ {stats.get('date_to', '')}", False, Pt(14), self.SECONDARY),
            (f"高リスク比率: {stats.get('high_risk_pct', 0):.2f}%", False, Pt(14), self.SECONDARY),
        ]
        for i, (text, bold, size, clr) in enumerate(extra_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.font.size = size
            p.font.bold = bold
            p.font.color.rgb = clr
            p.font.name = self.FONT_JP
            p.space_before = Pt(6)

    def _add_risk_chart_slide(self) -> None:
        """Slide 4: Risk distribution with doughnut chart."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, "リスク分布分析")
        self._add_slide_number(slide, 4)

        dist = self._risk_dist
        total = sum(dist.values()) or 1
        categories = ["高リスク (60点以上)", "中リスク (40-59点)",
                       "低リスク (20-39点)", "最小リスク (20点未満)"]
        values = [dist.get("high", 0), dist.get("medium", 0),
                  dist.get("low", 0), dist.get("minimal", 0)]

        if sum(values) > 0:
            # ドーナツチャート（左側）
            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series("件数", tuple(values))

            chart_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.DOUGHNUT,
                Inches(0.5), Inches(1.3), Inches(7), Inches(5.8),
                chart_data,
            )
            chart = chart_frame.chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False

            plot = chart.plots[0]
            series = plot.series[0]
            risk_colors = [self.COLOR_HIGH, self.COLOR_MEDIUM,
                           self.COLOR_LOW, self.COLOR_MINIMAL]
            for i, clr in enumerate(risk_colors):
                point = series.points[i]
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = clr

        # 右側テキスト: 数値詳細
        detail_box = slide.shapes.add_textbox(
            Inches(8), Inches(1.5), Inches(4.8), Inches(5.5),
        )
        tf = detail_box.text_frame
        tf.word_wrap = True

        detail_lines = [
            ("リスクレベル別内訳", True, Pt(18), self.PRIMARY),
            ("", False, Pt(6), self.SECONDARY),
        ]
        for cat, key, clr in [
            ("高リスク", "high", self.COLOR_HIGH),
            ("中リスク", "medium", self.COLOR_MEDIUM),
            ("低リスク", "low", self.COLOR_LOW),
            ("最小リスク", "minimal", self.COLOR_MINIMAL),
        ]:
            count = dist.get(key, 0)
            pct = count / total * 100
            detail_lines.append(
                (f"■ {cat}: {count:,}件 ({pct:.1f}%)", False, Pt(15), clr)
            )

        detail_lines.append(("", False, Pt(10), self.SECONDARY))
        detail_lines.append(
            ("高リスク仕訳には優先的な調査が必要です。", False, Pt(13), self.SECONDARY)
        )

        for i, (text, bold, size, clr) in enumerate(detail_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.font.size = size
            p.font.bold = bold
            p.font.color.rgb = clr
            p.font.name = self.FONT_JP
            p.space_before = Pt(4)

    def _add_findings_chart_slide(self) -> None:
        """Slide 5: Top findings with horizontal bar chart."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, "主要な発見事項")
        self._add_slide_number(slide, 5)

        findings = self._findings[:7]

        if findings:
            # 横棒グラフ（上位最大7件、逆順で表示）
            names = [f["rule_name"][:20] for f in reversed(findings)]
            counts = [f["count"] for f in reversed(findings)]

            chart_data = CategoryChartData()
            chart_data.categories = names
            chart_data.add_series("違反件数", tuple(counts))

            chart_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED,
                Inches(0.5), Inches(1.3), Inches(8.5), Inches(5.8),
                chart_data,
            )
            chart = chart_frame.chart
            chart.has_legend = False

            # バーの色を設定
            plot = chart.plots[0]
            series = plot.series[0]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = self.ACCENT

            # データラベル
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.size = Pt(10)
            data_labels.number_format = "#,##0"
            data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

            # 軸フォーマット
            chart.value_axis.has_title = False
            chart.category_axis.tick_labels.font.size = Pt(10)

            # 右側テキスト: トップ3の詳細
            detail_box = slide.shapes.add_textbox(
                Inches(9.3), Inches(1.5), Inches(3.7), Inches(5.5),
            )
            tf = detail_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "上位違反パターン"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.PRIMARY
            p.font.name = self.FONT_JP

            for i, f in enumerate(findings[:5], 1):
                p = tf.add_paragraph()
                p.text = f"{i}. {f['rule_name']}"
                p.font.size = Pt(12)
                p.font.color.rgb = self.SECONDARY
                p.font.name = self.FONT_JP
                p.space_before = Pt(8)

                p = tf.add_paragraph()
                p.text = f"   {f['count']:,}件検出"
                p.font.size = Pt(11)
                p.font.color.rgb = self.COLOR_HIGH if i <= 2 else self.SECONDARY
                p.font.name = self.FONT_JP
        else:
            no_data_box = slide.shapes.add_textbox(
                Inches(2), Inches(3), Inches(9), Inches(2),
            )
            tf = no_data_box.text_frame
            p = tf.paragraphs[0]
            p.text = "検出された重大な発見事項はありません。"
            p.font.size = Pt(20)
            p.font.color.rgb = self.COLOR_MINIMAL
            p.alignment = PP_ALIGN.CENTER
            p.font.name = self.FONT_JP

    def _add_benford_chart_slide(self) -> None:
        """Slide 6: Benford analysis with grouped column chart."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, "ベンフォードの法則分析")
        self._add_slide_number(slide, 6)

        digits = self._benford_digits
        summary = self._benford_summary

        if digits:
            # 集合縦棒グラフ: 実績 vs 理論値
            categories = [str(d["digit"]) for d in digits]
            actual_pcts = tuple(d["actual_pct"] for d in digits)
            expected_pcts = tuple(d["expected_pct"] for d in digits)

            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series("実績 (%)", actual_pcts)
            chart_data.add_series("理論値 (%)", expected_pcts)

            chart_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(0.5), Inches(1.3), Inches(8.5), Inches(5),
                chart_data,
            )
            chart = chart_frame.chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False

            # 系列の色を設定
            plot = chart.plots[0]
            actual_series = plot.series[0]
            actual_series.format.fill.solid()
            actual_series.format.fill.fore_color.rgb = self.ACCENT

            expected_series = plot.series[1]
            expected_series.format.fill.solid()
            expected_series.format.fill.fore_color.rgb = RGBColor(0xBD, 0xC3, 0xC7)

            # 値軸: パーセント表記
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.paragraphs[0].text = "%"
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)

            chart.category_axis.has_title = True
            chart.category_axis.axis_title.text_frame.paragraphs[0].text = "先頭桁"
            chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)

        # 右側: 適合度サマリー
        summary_box = slide.shapes.add_textbox(
            Inches(9.3), Inches(1.5), Inches(3.7), Inches(5.5),
        )
        tf = summary_box.text_frame
        tf.word_wrap = True

        conformity_color = self.COLOR_MINIMAL
        if summary.get("conformity") == "nonconforming":
            conformity_color = self.COLOR_HIGH
        elif summary.get("conformity") == "marginally_acceptable":
            conformity_color = self.COLOR_MEDIUM

        lines = [
            ("適合度分析結果", True, Pt(16), self.PRIMARY),
            ("", False, Pt(6), self.SECONDARY),
            (f"MAD: {summary.get('mad', 0):.4f}", False, Pt(15), self.SECONDARY),
            (f"評価: {summary.get('conformity_label', '未評価')}", True, Pt(15), conformity_color),
            ("", False, Pt(10), self.SECONDARY),
            ("【判定基準】", True, Pt(12), self.PRIMARY),
            ("0.006以下: 非常に良い適合", False, Pt(11), self.SECONDARY),
            ("0.012以下: 許容範囲", False, Pt(11), self.SECONDARY),
            ("0.015以下: 境界", False, Pt(11), self.SECONDARY),
            ("0.015超: 不適合（要調査）", False, Pt(11), self.SECONDARY),
        ]

        for i, (text, bold, size, clr) in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.font.size = size
            p.font.bold = bold
            p.font.color.rgb = clr
            p.font.name = self.FONT_JP
            p.space_before = Pt(3)

    def _add_trend_chart_slide(self) -> None:
        """Slide 7: Monthly trend with line chart."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, "月次トレンド分析")
        self._add_slide_number(slide, 7)

        trend = self._monthly_trend

        if trend:
            # 折れ線グラフ: 月次仕訳件数
            month_labels = [f"{t['period']}月" for t in trend]
            entry_counts = tuple(t["entry_count"] for t in trend)
            high_risk_counts = tuple(t["high_risk_count"] for t in trend)

            chart_data = CategoryChartData()
            chart_data.categories = month_labels
            chart_data.add_series("仕訳件数", entry_counts)
            chart_data.add_series("高リスク件数", high_risk_counts)

            chart_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE_MARKERS,
                Inches(0.5), Inches(1.3), Inches(9), Inches(5.5),
                chart_data,
            )
            chart = chart_frame.chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False

            # 系列色
            plot = chart.plots[0]
            entries_series = plot.series[0]
            entries_series.format.line.color.rgb = self.ACCENT
            entries_series.format.line.width = Pt(2.5)
            entries_series.smooth = False

            risk_series = plot.series[1]
            risk_series.format.line.color.rgb = self.COLOR_HIGH
            risk_series.format.line.width = Pt(2)
            risk_series.smooth = False

            # 右側テキスト: トレンド分析コメント
            max_month = max(trend, key=lambda t: t["entry_count"])
            max_risk_month = max(trend, key=lambda t: t["high_risk_count"])

            comment_box = slide.shapes.add_textbox(
                Inches(9.8), Inches(1.5), Inches(3.2), Inches(5.5),
            )
            tf = comment_box.text_frame
            tf.word_wrap = True

            comments = [
                ("トレンド概要", True, Pt(14), self.PRIMARY),
                ("", False, Pt(4), self.SECONDARY),
                (f"仕訳最多月: {max_month['period']}月", False, Pt(12), self.SECONDARY),
                (f"  ({max_month['entry_count']:,}件)", False, Pt(11), self.SECONDARY),
                ("", False, Pt(4), self.SECONDARY),
                (f"高リスク最多月: {max_risk_month['period']}月", False, Pt(12), self.COLOR_HIGH),
                (f"  ({max_risk_month['high_risk_count']:,}件)", False, Pt(11), self.COLOR_HIGH),
                ("", False, Pt(8), self.SECONDARY),
                ("期末月への仕訳集中は", False, Pt(11), self.SECONDARY),
                ("決算調整の可能性を", False, Pt(11), self.SECONDARY),
                ("示唆します。", False, Pt(11), self.SECONDARY),
            ]

            for i, (text, bold, size, clr) in enumerate(comments):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = text
                p.font.size = size
                p.font.bold = bold
                p.font.color.rgb = clr
                p.font.name = self.FONT_JP
        else:
            no_data_box = slide.shapes.add_textbox(
                Inches(2), Inches(3), Inches(9), Inches(2),
            )
            tf = no_data_box.text_frame
            p = tf.paragraphs[0]
            p.text = "月次トレンドデータがありません。"
            p.font.size = Pt(20)
            p.font.color.rgb = self.SECONDARY
            p.alignment = PP_ALIGN.CENTER
            p.font.name = self.FONT_JP

    def _add_recommendations_slide(self) -> None:
        """Slide 8: Data-driven recommendations."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, "推奨事項")
        self._add_slide_number(slide, 8)

        recs = self._generate_recommendations()

        content_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.5),
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "検証結果に基づく推奨事項:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.PRIMARY
        p.font.name = self.FONT_JP

        for i, rec in enumerate(recs, 1):
            # 番号行
            p = tf.add_paragraph()
            p.text = f"{i}. {rec['title']}"
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = rec.get("color", self.SECONDARY)
            p.font.name = self.FONT_JP
            p.space_before = Pt(14)

            # 詳細行
            p = tf.add_paragraph()
            p.text = f"   → {rec['detail']}"
            p.font.size = Pt(13)
            p.font.color.rgb = self.SECONDARY
            p.font.name = self.FONT_JP

    def _add_next_steps_slide(self) -> None:
        """Slide 9: Next steps with priority indicators."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, "次のステップ")
        self._add_slide_number(slide, 9)

        steps = self._generate_next_steps()

        content_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.5),
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "フォローアップ項目:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.PRIMARY
        p.font.name = self.FONT_JP

        for step in steps:
            priority = step["priority"]
            if priority == "高":
                priority_color = self.COLOR_HIGH
            elif priority == "中":
                priority_color = self.COLOR_MEDIUM
            else:
                priority_color = self.COLOR_LOW

            p = tf.add_paragraph()
            p.text = f"[{priority}] {step['task']}"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = priority_color
            p.font.name = self.FONT_JP
            p.space_before = Pt(12)

            p = tf.add_paragraph()
            p.text = f"      期限: {step['deadline']}　担当: [          ]"
            p.font.size = Pt(12)
            p.font.color.rgb = self.SECONDARY
            p.font.name = self.FONT_JP

    def _add_appendix_slide(self) -> None:
        """Slide 10: Appendix - methodology."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, "付録：検証手法")
        self._add_slide_number(slide, 10)

        # 左列: ルールベース
        left_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.5), Inches(3.8), Inches(5.5),
        )
        tf = left_box.text_frame
        tf.word_wrap = True

        left_lines = [
            ("ルールベース検証", True, Pt(16)),
            ("", False, Pt(4)),
            ("58種類の監査ルールを適用", False, Pt(13)),
            ("", False, Pt(4)),
            ("カテゴリ:", True, Pt(13)),
            ("• 金額チェック", False, Pt(12)),
            ("• 時間帯チェック", False, Pt(12)),
            ("• 勘定科目チェック", False, Pt(12)),
            ("• 承認チェック", False, Pt(12)),
            ("• ML異常検出", False, Pt(12)),
            ("• ベンフォード分析", False, Pt(12)),
        ]
        for i, (text, bold, size) in enumerate(left_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.font.size = size
            p.font.bold = bold
            p.font.color.rgb = self.PRIMARY if bold else self.SECONDARY
            p.font.name = self.FONT_JP

        # 中央列: ML手法
        mid_box = slide.shapes.add_textbox(
            Inches(5), Inches(1.5), Inches(3.8), Inches(5.5),
        )
        tf = mid_box.text_frame
        tf.word_wrap = True

        mid_lines = [
            ("機械学習異常検出", True, Pt(16)),
            ("", False, Pt(4)),
            ("5手法のアンサンブル:", True, Pt(13)),
            ("", False, Pt(4)),
            ("• Isolation Forest", False, Pt(12)),
            ("• Local Outlier Factor", False, Pt(12)),
            ("• One-Class SVM", False, Pt(12)),
            ("• Autoencoder", False, Pt(12)),
            ("• アンサンブル投票", False, Pt(12)),
        ]
        for i, (text, bold, size) in enumerate(mid_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.font.size = size
            p.font.bold = bold
            p.font.color.rgb = self.PRIMARY if bold else self.SECONDARY
            p.font.name = self.FONT_JP

        # 右列: 統計分析
        right_box = slide.shapes.add_textbox(
            Inches(9.2), Inches(1.5), Inches(3.8), Inches(5.5),
        )
        tf = right_box.text_frame
        tf.word_wrap = True

        right_lines = [
            ("統計分析", True, Pt(16)),
            ("", False, Pt(4)),
            ("ベンフォードの法則:", True, Pt(13)),
            ("• 第1桁分布分析", False, Pt(12)),
            ("• MAD適合度評価", False, Pt(12)),
            ("", False, Pt(8)),
            ("リスクスコアリング:", True, Pt(13)),
            ("• ルール違反 (60%)", False, Pt(12)),
            ("• ML異常 (40%)", False, Pt(12)),
            ("• 5段階分類", False, Pt(12)),
        ]
        for i, (text, bold, size) in enumerate(right_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.font.size = size
            p.font.bold = bold
            p.font.color.rgb = self.PRIMARY if bold else self.SECONDARY
            p.font.name = self.FONT_JP

    # ========== データ駆動コンテンツ生成 ==========

    def _generate_recommendations(self) -> list[dict[str, Any]]:
        """Generate data-driven recommendations."""
        recs = []
        stats = self._stats
        benford = self._benford_summary
        findings = self._findings
        high_risk_pct = stats.get("high_risk_pct", 0)
        high_risk_count = stats.get("high_risk_count", 0)

        if high_risk_pct > 5:
            recs.append({
                "title": f"高リスク仕訳の緊急調査（{high_risk_count:,}件）",
                "detail": f"高リスク比率が{high_risk_pct:.1f}%と高水準です。即座に個別調査を実施してください。",
                "color": self.COLOR_HIGH,
            })
        elif high_risk_count > 0:
            recs.append({
                "title": f"高リスク仕訳の調査（{high_risk_count:,}件）",
                "detail": "担当者へのヒアリングと証憑確認を実施してください。",
                "color": self.COLOR_MEDIUM,
            })

        if benford.get("conformity") in ("nonconforming", "marginally_acceptable"):
            mad = benford.get("mad", 0)
            recs.append({
                "title": "ベンフォード分析の追加調査",
                "detail": f"MAD={mad:.4f}で{benford.get('conformity_label', '')}。データ操作の可能性を調査してください。",
                "color": self.COLOR_HIGH if benford.get("conformity") == "nonconforming" else self.COLOR_MEDIUM,
            })

        if findings:
            top = findings[0]
            recs.append({
                "title": "最多違反パターンの原因分析",
                "detail": f"「{top['rule_name']}」({top['count']:,}件)について根本原因を究明してください。",
                "color": self.SECONDARY,
            })

        recs.append({
            "title": "内部統制の改善提案",
            "detail": "検出された問題点に対する改善策を策定してください。",
            "color": self.SECONDARY,
        })
        recs.append({
            "title": "定期的モニタリングの継続",
            "detail": "AI監査ツールによる継続的なモニタリング体制を構築してください。",
            "color": self.SECONDARY,
        })

        return recs

    def _generate_next_steps(self) -> list[dict[str, str]]:
        """Generate data-driven next steps."""
        steps = []
        stats = self._stats
        findings = self._findings
        benford = self._benford_summary

        high_risk = stats.get("high_risk_count", 0)
        if high_risk > 0:
            steps.append({
                "priority": "高",
                "task": f"高リスク仕訳 {high_risk:,}件の個別調査",
                "deadline": "2週間以内",
            })

        if benford.get("conformity") in ("nonconforming", "marginally_acceptable"):
            steps.append({
                "priority": "高",
                "task": "ベンフォード分析結果の深掘り調査",
                "deadline": "1ヶ月以内",
            })

        if findings:
            steps.append({
                "priority": "中",
                "task": f"違反パターン上位{min(len(findings), 5)}件の原因分析",
                "deadline": "1ヶ月以内",
            })

        steps.append({
            "priority": "中",
            "task": "内部統制改善提案の策定",
            "deadline": "四半期内",
        })
        steps.append({
            "priority": "低",
            "task": "経営陣への報告書提出",
            "deadline": "1ヶ月以内",
        })
        steps.append({
            "priority": "低",
            "task": "是正措置の実施確認",
            "deadline": "四半期内",
        })

        return steps


class PDFReportGenerator:
    """PDF report generator with embedded charts."""

    # matplotlib用カラー定義
    COLORS = {
        "primary": "#1A365D",
        "accent": "#2980B9",
        "high": "#E74C3C",
        "medium": "#F39C12",
        "low": "#2980B9",
        "minimal": "#27AE60",
        "gray": "#BDC3C7",
    }

    def __init__(self, config: ReportConfig):
        """Initialize PDF generator."""
        self.config = config
        self.db = DuckDBManager()
        self.styles = getSampleStyleSheet()
        self._setup_styles()

        # データキャッシュ
        self._stats: dict[str, Any] = {}
        self._risk_dist: dict[str, int] = {}
        self._findings: list[dict[str, Any]] = []
        self._benford_digits: list[dict[str, Any]] = []
        self._benford_summary: dict[str, Any] = {}
        self._monthly_trend: list[dict[str, Any]] = []

    def _setup_styles(self) -> None:
        """Setup custom styles."""
        self.styles.add(
            ParagraphStyle(
                name="JapaneseTitle",
                fontSize=24,
                leading=30,
                alignment=1,
                spaceAfter=20,
            )
        )
        self.styles.add(
            ParagraphStyle(
                name="JapaneseHeading",
                fontSize=16,
                leading=20,
                spaceBefore=15,
                spaceAfter=10,
            )
        )
        self.styles.add(
            ParagraphStyle(
                name="JapaneseBody",
                fontSize=10,
                leading=14,
            )
        )

    def generate(self) -> bytes:
        """Generate PDF report with charts.

        Returns:
            PDF file as bytes.
        """
        self._prefetch_data()

        buffer = io.BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=A4,
            rightMargin=20 * mm,
            leftMargin=20 * mm,
            topMargin=20 * mm,
            bottomMargin=20 * mm,
        )

        story = []

        # Title
        story.append(Paragraph(self.config.report_title, self.styles["JapaneseTitle"]))
        story.append(
            Paragraph(
                f"{self.config.fiscal_year}年度 {self.config.company_name}",
                self.styles["Normal"],
            )
        )
        story.append(
            Paragraph(
                f"作成日: {datetime.now().strftime('%Y年%m月%d日')}　作成者: {self.config.prepared_by}",
                self.styles["Normal"],
            )
        )
        story.append(Spacer(1, 30))

        # 1. Executive Summary
        story.append(
            Paragraph("1. エグゼクティブサマリー", self.styles["JapaneseHeading"])
        )
        stats = self._stats
        summary_data = [
            ["項目", "値"],
            ["検証対象仕訳件数", f"{stats.get('total_entries', 0):,}件"],
            ["仕訳帳票数", f"{stats.get('total_journals', 0):,}件"],
            ["総取引金額", f"¥{stats.get('total_amount', 0):,.0f}"],
            ["高リスク仕訳", f"{stats.get('high_risk_count', 0):,}件 ({stats.get('high_risk_pct', 0):.1f}%)"],
            ["平均リスクスコア", f"{stats.get('avg_risk_score', 0):.1f}点"],
            ["対象期間", f"{stats.get('date_from', '')} ～ {stats.get('date_to', '')}"],
        ]
        story.append(self._create_table(summary_data, header=True))
        story.append(Spacer(1, 20))

        # 2. Risk Distribution with chart
        story.append(Paragraph("2. リスク分布", self.styles["JapaneseHeading"]))
        dist = self._risk_dist
        dist_data = [
            ["リスクレベル", "件数", "割合"],
            ["高リスク (60点以上)", f"{dist.get('high', 0):,}", f"{self._pct(dist.get('high', 0))}%"],
            ["中リスク (40-59点)", f"{dist.get('medium', 0):,}", f"{self._pct(dist.get('medium', 0))}%"],
            ["低リスク (20-39点)", f"{dist.get('low', 0):,}", f"{self._pct(dist.get('low', 0))}%"],
            ["最小リスク (20点未満)", f"{dist.get('minimal', 0):,}", f"{self._pct(dist.get('minimal', 0))}%"],
        ]
        story.append(self._create_table(dist_data, header=True, col_widths=[60, 40, 30]))
        story.append(Spacer(1, 10))

        # リスク分布円グラフ
        risk_chart = self._create_risk_pie_chart()
        if risk_chart:
            story.append(Image(risk_chart, width=130 * mm, height=80 * mm))
        story.append(Spacer(1, 20))

        # 3. Benford Analysis with chart
        story.append(
            Paragraph("3. ベンフォードの法則分析", self.styles["JapaneseHeading"])
        )
        summary = self._benford_summary
        story.append(
            Paragraph(
                f"MAD (平均絶対偏差): {summary.get('mad', 0):.4f}　"
                f"評価: {summary.get('conformity_label', '未評価')}",
                self.styles["JapaneseBody"],
            )
        )
        story.append(Spacer(1, 10))

        benford_chart = self._create_benford_bar_chart()
        if benford_chart:
            story.append(Image(benford_chart, width=150 * mm, height=80 * mm))
        story.append(Spacer(1, 20))

        # 4. Monthly Trend with chart
        story.append(PageBreak())
        story.append(
            Paragraph("4. 月次トレンド分析", self.styles["JapaneseHeading"])
        )
        trend_chart = self._create_trend_line_chart()
        if trend_chart:
            story.append(Image(trend_chart, width=150 * mm, height=80 * mm))
        story.append(Spacer(1, 20))

        # 5. Top Findings
        story.append(Paragraph("5. 主要な発見事項", self.styles["JapaneseHeading"]))
        findings = self._findings
        if findings:
            findings_data = [["順位", "ルール名", "違反件数"]]
            for i, f in enumerate(findings[:10], 1):
                findings_data.append([str(i), f["rule_name"], f"{f['count']:,}"])
            story.append(self._create_table(findings_data, header=True, col_widths=[15, 95, 30]))
        else:
            story.append(
                Paragraph(
                    "検出された重大な発見事項はありません。",
                    self.styles["JapaneseBody"],
                )
            )
        story.append(Spacer(1, 20))

        # 6. Recommendations
        story.append(Paragraph("6. 推奨事項", self.styles["JapaneseHeading"]))
        recs = self._generate_recommendations()
        for i, rec in enumerate(recs, 1):
            story.append(
                Paragraph(f"{i}. {rec}", self.styles["JapaneseBody"])
            )
            story.append(Spacer(1, 5))

        # Build PDF
        doc.build(story)
        buffer.seek(0)
        return buffer.getvalue()

    # ========== データ取得 ==========

    def _prefetch_data(self) -> None:
        """Pre-fetch all data for report generation."""
        self._stats = self._query_summary_stats()
        self._risk_dist = self._query_risk_distribution()
        self._findings = self._query_top_findings()
        self._benford_digits, self._benford_summary = self._query_benford_data()
        self._monthly_trend = self._query_monthly_trend()

    def _query_summary_stats(self) -> dict[str, Any]:
        """Get summary statistics."""
        query = """
            SELECT
                COUNT(*) as total_entries,
                COUNT(DISTINCT journal_id) as total_journals,
                COALESCE(SUM(ABS(amount)), 0) as total_amount,
                MIN(effective_date) as date_from,
                MAX(effective_date) as date_to,
                SUM(CASE WHEN risk_score >= 60 THEN 1 ELSE 0 END) as high_risk_count,
                AVG(CASE WHEN risk_score > 0 THEN risk_score END) as avg_risk_score
            FROM journal_entries
            WHERE fiscal_year = ?
        """
        result = self.db.execute(query, [self.config.fiscal_year])
        if result:
            row = result[0]
            total = row[0] or 1
            return {
                "total_entries": row[0] or 0,
                "total_journals": row[1] or 0,
                "total_amount": row[2] or 0,
                "date_from": str(row[3]) if row[3] else "",
                "date_to": str(row[4]) if row[4] else "",
                "high_risk_count": row[5] or 0,
                "high_risk_pct": (row[5] or 0) / total * 100,
                "avg_risk_score": row[6] or 0,
            }
        return {}

    def _query_risk_distribution(self) -> dict[str, int]:
        """Get risk distribution."""
        query = """
            SELECT
                CASE
                    WHEN risk_score >= 60 THEN 'high'
                    WHEN risk_score >= 40 THEN 'medium'
                    WHEN risk_score >= 20 THEN 'low'
                    ELSE 'minimal'
                END as level,
                COUNT(*) as count
            FROM journal_entries
            WHERE fiscal_year = ?
            GROUP BY level
        """
        result = self.db.execute(query, [self.config.fiscal_year])
        dist = {"high": 0, "medium": 0, "low": 0, "minimal": 0}
        for row in result:
            dist[row[0]] = row[1]
        return dist

    def _query_top_findings(self) -> list[dict[str, Any]]:
        """Get top findings."""
        query = """
            SELECT rule_name, COUNT(*) as count
            FROM rule_violations rv
            JOIN journal_entries je ON rv.gl_detail_id = je.gl_detail_id
            WHERE je.fiscal_year = ?
            GROUP BY rule_name
            ORDER BY count DESC
            LIMIT 10
        """
        result = self.db.execute(query, [self.config.fiscal_year])
        return [{"rule_name": row[0], "count": row[1]} for row in result]

    def _query_benford_data(self) -> tuple[list[dict], dict]:
        """Get Benford analysis with per-digit data."""
        query = """
            SELECT
                CAST(SUBSTR(CAST(ABS(CAST(amount AS BIGINT)) AS VARCHAR), 1, 1) AS INTEGER) as digit,
                COUNT(*) as count
            FROM journal_entries
            WHERE fiscal_year = ?
                AND ABS(amount) >= 10
            GROUP BY digit
            HAVING digit BETWEEN 1 AND 9
            ORDER BY digit
        """
        result = self.db.execute(query, [self.config.fiscal_year])

        total = sum(row[1] for row in result) or 1
        digits = []
        mad = 0
        for digit_num in range(1, 10):
            count = 0
            for row in result:
                if row[0] == digit_num:
                    count = row[1]
                    break
            actual = count / total
            expected = BENFORD_EXPECTED.get(digit_num, 0)
            mad += abs(actual - expected)
            digits.append({
                "digit": digit_num,
                "actual_pct": round(actual * 100, 2),
                "expected_pct": round(expected * 100, 2),
            })
        mad = mad / 9

        if mad <= 0.006:
            label, conformity = "非常に良い適合", "close"
        elif mad <= 0.012:
            label, conformity = "許容範囲", "acceptable"
        elif mad <= 0.015:
            label, conformity = "境界", "marginally_acceptable"
        else:
            label, conformity = "不適合（要調査）", "nonconforming"

        return digits, {"mad": mad, "conformity": conformity, "conformity_label": label}

    def _query_monthly_trend(self) -> list[dict[str, Any]]:
        """Get monthly trend data."""
        query = """
            SELECT
                accounting_period as period,
                COUNT(*) as entry_count,
                COALESCE(SUM(ABS(amount)), 0) as total_amount,
                SUM(CASE WHEN risk_score >= 60 THEN 1 ELSE 0 END) as high_risk_count
            FROM journal_entries
            WHERE fiscal_year = ?
            GROUP BY accounting_period
            ORDER BY accounting_period
        """
        result = self.db.execute(query, [self.config.fiscal_year])
        return [
            {"period": row[0], "entry_count": row[1],
             "total_amount": row[2], "high_risk_count": row[3]}
            for row in result
        ]

    # ========== チャート生成（matplotlib → BytesIO） ==========

    def _create_risk_pie_chart(self) -> io.BytesIO | None:
        """Create risk distribution pie chart."""
        dist = self._risk_dist
        labels = ["高リスク", "中リスク", "低リスク", "最小リスク"]
        values = [dist.get("high", 0), dist.get("medium", 0),
                  dist.get("low", 0), dist.get("minimal", 0)]
        colors = [self.COLORS["high"], self.COLORS["medium"],
                  self.COLORS["low"], self.COLORS["minimal"]]

        if sum(values) == 0:
            return None

        fig, ax = plt.subplots(figsize=(8, 5))
        wedges, texts, autotexts = ax.pie(
            values, labels=labels, colors=colors, autopct="%1.1f%%",
            startangle=90, pctdistance=0.75,
            wedgeprops={"width": 0.4, "edgecolor": "white", "linewidth": 2},
        )
        for t in texts:
            t.set_fontsize(10)
        for t in autotexts:
            t.set_fontsize(9)
            t.set_color("white")
            t.set_weight("bold")

        ax.set_title("リスクレベル別分布", fontsize=14, fontweight="bold", pad=15)
        fig.tight_layout()

        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                    facecolor="white", edgecolor="none")
        buf.seek(0)
        plt.close(fig)
        return buf

    def _create_benford_bar_chart(self) -> io.BytesIO | None:
        """Create Benford analysis bar chart."""
        digits = self._benford_digits
        if not digits:
            return None

        x = [str(d["digit"]) for d in digits]
        actual = [d["actual_pct"] for d in digits]
        expected = [d["expected_pct"] for d in digits]

        fig, ax = plt.subplots(figsize=(9, 5))
        bar_width = 0.35
        x_pos = range(len(x))

        ax.bar(
            [p - bar_width / 2 for p in x_pos], actual, bar_width,
            label="実績 (%)", color=self.COLORS["accent"], edgecolor="white",
        )
        ax.bar(
            [p + bar_width / 2 for p in x_pos], expected, bar_width,
            label="理論値 (%)", color=self.COLORS["gray"], edgecolor="white",
        )

        ax.set_xlabel("先頭桁", fontsize=11)
        ax.set_ylabel("%", fontsize=11)
        ax.set_title("ベンフォードの法則: 実績 vs 理論値", fontsize=14, fontweight="bold")
        ax.set_xticks(list(x_pos))
        ax.set_xticklabels(x)
        ax.legend(fontsize=10)
        ax.grid(axis="y", alpha=0.3)
        fig.tight_layout()

        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                    facecolor="white", edgecolor="none")
        buf.seek(0)
        plt.close(fig)
        return buf

    def _create_trend_line_chart(self) -> io.BytesIO | None:
        """Create monthly trend line chart."""
        trend = self._monthly_trend
        if not trend:
            return None

        months = [f"{t['period']}月" for t in trend]
        counts = [t["entry_count"] for t in trend]
        high_risk = [t["high_risk_count"] for t in trend]

        fig, ax1 = plt.subplots(figsize=(9, 5))

        ax1.plot(months, counts, "o-", color=self.COLORS["accent"],
                 linewidth=2, markersize=6, label="仕訳件数")
        ax1.set_xlabel("月", fontsize=11)
        ax1.set_ylabel("仕訳件数", color=self.COLORS["accent"], fontsize=11)
        ax1.tick_params(axis="y", labelcolor=self.COLORS["accent"])
        ax1.yaxis.set_major_formatter(mticker.FuncFormatter(
            lambda x, _p: f"{x:,.0f}"
        ))

        # 高リスク件数（第2軸）
        ax2 = ax1.twinx()
        ax2.plot(months, high_risk, "s--", color=self.COLORS["high"],
                 linewidth=1.5, markersize=5, label="高リスク件数")
        ax2.set_ylabel("高リスク件数", color=self.COLORS["high"], fontsize=11)
        ax2.tick_params(axis="y", labelcolor=self.COLORS["high"])

        # 凡例の統合
        lines1, labels1 = ax1.get_legend_handles_labels()
        lines2, labels2 = ax2.get_legend_handles_labels()
        ax1.legend(lines1 + lines2, labels1 + labels2, loc="upper left", fontsize=10)

        ax1.set_title("月次トレンド: 仕訳件数 & 高リスク件数",
                       fontsize=14, fontweight="bold")
        ax1.grid(axis="y", alpha=0.3)
        plt.xticks(rotation=45 if len(months) > 6 else 0)
        fig.tight_layout()

        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                    facecolor="white", edgecolor="none")
        buf.seek(0)
        plt.close(fig)
        return buf

    # ========== ヘルパー ==========

    def _pct(self, count: int) -> str:
        """Calculate percentage string."""
        total = sum(self._risk_dist.values()) or 1
        return f"{count / total * 100:.1f}"

    def _create_table(
        self, data: list[list[str]], header: bool = False,
        col_widths: list[int] | None = None,
    ) -> Table:
        """Create formatted table."""
        widths = [w * mm for w in col_widths] if col_widths else [100 * mm, 50 * mm]

        table = Table(data, colWidths=widths)
        style = [
            ("GRID", (0, 0), (-1, -1), 0.5, rl_colors.grey),
            ("FONTSIZE", (0, 0), (-1, -1), 10),
            ("PADDING", (0, 0), (-1, -1), 5),
        ]
        if header:
            style.extend(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), rl_colors.HexColor("#1A365D")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), rl_colors.white),
                    ("FONTSIZE", (0, 0), (-1, 0), 11),
                ]
            )
        # 交互行の背景色
        for i in range(1, len(data)):
            if i % 2 == 0:
                style.append(
                    ("BACKGROUND", (0, i), (-1, i), rl_colors.HexColor("#F0F2F5"))
                )
        table.setStyle(TableStyle(style))
        return table

    def _generate_recommendations(self) -> list[str]:
        """Generate data-driven recommendation strings for PDF."""
        recs = []
        stats = self._stats
        benford = self._benford_summary
        findings = self._findings
        high_risk_pct = stats.get("high_risk_pct", 0)
        high_risk_count = stats.get("high_risk_count", 0)

        if high_risk_pct > 5:
            recs.append(
                f"【緊急】高リスク仕訳が{high_risk_count:,}件（{high_risk_pct:.1f}%）検出されました。"
                f"即座に個別調査を実施してください。"
            )
        elif high_risk_count > 0:
            recs.append(
                f"高リスク仕訳{high_risk_count:,}件について、担当者へのヒアリングと証憑確認を実施してください。"
            )

        if benford.get("conformity") in ("nonconforming", "marginally_acceptable"):
            recs.append(
                f"ベンフォード分析で{benford.get('conformity_label', '')}が検出されました"
                f"（MAD: {benford.get('mad', 0):.4f}）。データ操作の可能性を調査してください。"
            )

        if findings:
            top = findings[0]
            recs.append(
                f"最多違反パターン「{top['rule_name']}」（{top['count']:,}件）について原因分析を実施してください。"
            )

        recs.append("内部統制の有効性を評価し、改善計画を策定してください。")
        recs.append("AI監査ツールによる継続的モニタリング体制の構築を推奨します。")

        return recs
